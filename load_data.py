import os
import PyPDF2
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_ppt(fp):
    import platform
    os_type = platform.system()
    if os_type == "Windows":
        import win32com.client as wc
        powerpoint = wc.Dispatch("PowerPoint.Application")
        wc.gencache.EnsureDispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        ppt = powerpoint.Presentations.Open(fp)
        ppt.SaveAs(fp + "x")
        powerpoint.Quit()
    elif os_type == "Linux":
        import aspose.slides as slides
        with slides.Presentation(fp) as presentation:
            presentation.save(fp + "x", slides.export.SaveFormat.PPTX)
    read_pptx(fp + "x")


# 打开文件，可打开各个类型的文字文件
def read_pptx(file_path):
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        # print(f"Slide {slide_num}:")

        for shape in slide.shapes:
            # print(f"  Shape: {shape.name}, Type: {shape.shape_type}")
            # 检查形状是否有 text_frame 属性
            if hasattr(shape, "text_frame"):
                # 如果形状有 text_frame，则遍历其中的段落并打印文本
                for paragraph in shape.text_frame.paragraphs:
                    # print(f"    Paragraph: {paragraph.text}")
                    print(f"{paragraph.text}")
            # if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            #     if hasattr(shape, "text_frame"):
            #         for paragraph in shape.text_frame.paragraphs:
            #             print(f"  Paragraph: {paragraph.text}")

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # print(f"  Image found: {shape.image.blob}")
                # print(f"  Image found in page {slide_num}")
                print(f"Image found in page {file_path}-{slide_num}")

            # 你可以根据需要添加更多的形状类型处理
def open_file(file_path):
    """
    根据文件路径打开文件，并打印文件内容
    :param file_path: 文件路径
    :return: None
    """
    # txt文件
    if file_path.endswith('.txt'):
        print("读取txt文件...")
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
            print(content)
    # pdf文件
    elif file_path.endswith('.pdf'):
        print("读取pdf文件...")
        # 使用PyPDF2库打开PDF文件
        pdf_file = open(file_path, 'rb')
        pdf_reader = PyPDF2.PdfReader(file_path)
        for page in pdf_reader.pages:
            print(page.extract_text())
    # docx文档
    elif file_path.endswith('.docx'):
        print("读取docx文件...")
        doc = Document(file_path)
        # 遍历所有段落并打印其内容
        for para in doc.paragraphs:
            print(para.text)
    # pptx文件
    elif file_path.endswith('.pptx'):
        print("读取pptx文件...")
        read_pptx(file_path)
    # ppt文件
    elif file_path.endswith('.ppt'):
        print("读取ppt文件...")
        read_ppt(file_path)


# with open(file_path, 'r', encoding='utf-8') as file:
#     content = file.read()
#     print(content)

# 遍历文件夹中的所有文件
def list_files(folder_path):
    """
    遍历文件夹中的所有文件
    :param folder_path: 文件夹路径
    :return: None
    """
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if os.path.isfile(file_path):
            print(file_path)


if __name__ == '__main__':
    path = input("请输入文件路径：")
    # 判断是不是有效路径
    if os.path.exists(path):
        # 如果路径是个文件夹
        if os.path.isdir(path):
            list_files(path)
        # 如果路径是txt、docx、pdf文件
        elif path.endswith('.txt') or path.endswith('.pdf') or path.endswith('.docx') or path.endswith('.ppt') or path.endswith('.pptx'):
            open_file(path)
    else:
        print("文件路径不存在！进入测试阶段...")
        # 文件夹、文件路径 .txt、.pdf、.docx
        folder_path = "../build_data"
        chonce = input("按键‘1’：打开文件；按键‘2’：打开文件夹；任意按键：退出程序")

        if chonce == '1':
            file = input("按键‘1’：选择txt文件；按键‘2’：选择docx文件；按键‘3’：选择pdf文件；任意按键：退出程序")
            if file == '1':
                file_path = "./test.txt"
            elif file == '2':
                file_path = "./test.docx"
            elif file == '3':
                file_path = "./test.pdf"
            else:
                exit("程序已退出！")
            open_file(file_path)
        elif chonce == '2':
            list_files(folder_path)
        else:

            exit("程序已退出！")
